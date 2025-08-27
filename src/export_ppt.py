from __future__ import annotations
from typing import Dict
from pathlib import Path
from datetime import datetime
import tempfile
import plotly.io as pio
from pptx import Presentation
from pptx.util import Inches

TEMPLATES_DIR = Path("templates")
DEFAULT_TEMPLATE = TEMPLATES_DIR / "ppt_template.pptx"

def _ensure_template() -> Path:
    """
    Devuelve la plantilla si existe; si no, crea una muy básica automáticamente.
    """
    if DEFAULT_TEMPLATE.exists():
        return DEFAULT_TEMPLATE
    TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    tmp = TEMPLATES_DIR / "_auto_template.pptx"
    prs.save(tmp)
    return tmp

def export_ppt(institution: str, kpis: dict, figs: Dict[str, "plotly.graph_objs._figure.Figure"], comments: str) -> str:
    """
    Exporta un PPT con portada, KPIs, gráficos e interpretación.
    Devuelve la ruta absoluta del archivo generado (útil para st.download_button).
    """
    prs = Presentation(_ensure_template())

    # Portada
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"Reporte - {institution}"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = f"Fecha: {datetime.now():%Y-%m-%d}\nGenerado automáticamente"
    else:
        # Si la plantilla no tiene placeholder de subtítulo, añadimos un textbox
        tf = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(9), Inches(1.2)).text_frame
        tf.word_wrap = True
        tf.text = f"Fecha: {datetime.now():%Y-%m-%d}\nGenerado automáticamente"

    # KPIs
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "KPIs"
    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(9), Inches(3)).text_frame
    body.word_wrap = True
    body.text = (
        f"Usuarios: {kpis.get('usuarios_total', 0)}\n"
        f"Activos 90d: {kpis.get('activos_90d_pct', 0):.1f}%\n"
        f"Experiencia mediana: {kpis.get('exp_mediana_anos', 0):.1f} años\n"
        f"Salario mediano: S/ {kpis.get('sal_mediana', 0):.0f}"
    )

    # Gráficos
    for title, fig in (figs or {}).items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        img_path = None
        try:
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                # Requiere 'kaleido' instalado en el entorno
                pio.write_image(fig, tmp.name, width=1280, height=720, scale=2)
                img_path = tmp.name
        except Exception:
            # Si no hay kaleido o falla el render, dejamos una nota
            tf = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(9), Inches(3)).text_frame
            tf.text = "(No se pudo exportar el gráfico como imagen. Instala 'kaleido' para incluir gráficos en el PPT.)"
        if img_path:
            slide.shapes.add_picture(img_path, Inches(0.7), Inches(1.5), Inches(9), Inches(5))

    # Comentarios / narrativa
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Comentarios"
    tf = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(9), Inches(4.5)).text_frame
    tf.word_wrap = True
    tf.text = comments or "(Sin comentarios)"

    out = Path(f"Reporte_{institution}.pptx")
    prs.save(out)
    return str(out.absolute())
